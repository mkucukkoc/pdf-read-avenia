from __future__ import annotations

import io
import os
import re
import subprocess
from dataclasses import dataclass
from typing import Literal, List, Dict

import logging
from pypdf import PdfReader
from pdf2image import convert_from_bytes
from PIL import Image
import pytesseract
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)
logger.debug("[doc_text] module imported")


# ---------------- Dosya tespiti -----------------
def detect_file_type(filename: str, content_type: str | None) -> Literal["pdf", "docx", "pptx", "ppt", "unknown"]:
    """Dosya uzantısı ve MIME tipine göre belirleme."""
    logger.debug("[detect_file_type] filename=%s content_type=%s", filename, content_type)
    name = (filename or "").lower()
    if name.endswith(".pdf") or content_type == "application/pdf":
        logger.info("[detect_file_type] -> pdf")
        return "pdf"
    if name.endswith(".docx") or content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        logger.info("[detect_file_type] -> docx")
        return "docx"
    if name.endswith(".pptx") or content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        logger.info("[detect_file_type] -> pptx")
        return "pptx"
    if name.endswith(".ppt"):
        logger.info("[detect_file_type] -> ppt (legacy)")
        return "ppt"
    logger.info("[detect_file_type] -> unknown")
    return "unknown"


# ---------------- PDF -----------------
def extract_text_from_pdf_bytes(pdf_bytes: bytes) -> str:
    logger.info("[extract_text_from_pdf_bytes] start bytes_len=%s", len(pdf_bytes or b""))
    reader = PdfReader(io.BytesIO(pdf_bytes))
    try:
        page_count = len(reader.pages)
    except Exception:
        page_count = 0
    logger.debug("[extract_text_from_pdf_bytes] page_count=%s", page_count)
    texts: List[str] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        logger.debug("[extract_text_from_pdf_bytes] page=%s text_len=%s preview=%r", idx, len(text), text[:200])
        texts.append(text)
    joined = "\n".join(texts)
    logger.info("[extract_text_from_pdf_bytes] end total_text_len=%s", len(joined))
    return joined


def extract_text_via_ocr(pdf_bytes: bytes, max_pages: int | None = None) -> str:
    logger.info("[extract_text_via_ocr] start bytes_len=%s max_pages=%s", len(pdf_bytes or b""), max_pages)
    images = convert_from_bytes(pdf_bytes, first_page=1, last_page=max_pages)
    logger.debug("[extract_text_via_ocr] images_count=%s", len(images))
    ocr_texts = []
    for i, image in enumerate(images, start=1):
        txt = pytesseract.image_to_string(image)
        logger.debug("[extract_text_via_ocr] page_image=%s ocr_text_len=%s preview=%r", i, len(txt), txt[:200])
        ocr_texts.append(txt)
    joined = "\n".join(ocr_texts)
    logger.info("[extract_text_via_ocr] end total_ocr_text_len=%s", len(joined))
    return joined


def split_pdf_by_pages(pdf_bytes: bytes, max_chars: int) -> List[Dict[str, str]]:
    logger.info("[split_pdf_by_pages] start bytes_len=%s max_chars=%s", len(pdf_bytes or b""), max_chars)
    reader = PdfReader(io.BytesIO(pdf_bytes))
    try:
        page_count = len(reader.pages)
    except Exception:
        page_count = 0
    logger.debug("[split_pdf_by_pages] page_count=%s", page_count)
    chunks: List[Dict[str, str]] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = normalize_text(text)
        parts = split_text_by_size(text, max_chars)
        logger.debug("[split_pdf_by_pages] page=%s normalized_len=%s parts=%s", i, len(text), len(parts))
        for part in parts:
            chunks.append({"source": f"page:{i}", "text": part})
    logger.info("[split_pdf_by_pages] end chunks_count=%s", len(chunks))
    return chunks


def extract_images_from_pdf_bytes(pdf_bytes: bytes) -> List[Dict[str, Image.Image]]:
    """PDF içindeki görüntüleri çıkarır."""
    logger.info("[extract_images_from_pdf_bytes] start bytes_len=%s", len(pdf_bytes or b""))
    images: List[Dict[str, Image.Image]] = []
    reader = PdfReader(io.BytesIO(pdf_bytes))
    for p_idx, page in enumerate(reader.pages, start=1):
        try:
            page_images = getattr(page, "images", [])
            logger.debug("[extract_images_from_pdf_bytes] page=%s image_count=%s", p_idx, len(page_images))
            for i_idx, img in enumerate(page_images, start=1):
                data = getattr(img, "data", b"")
                logger.debug(
                    "[extract_images_from_pdf_bytes] page=%s img=%s data_len=%s",
                    p_idx,
                    i_idx,
                    len(data or b""),
                )
                if not data:
                    continue
                try:
                    pil_img = Image.open(io.BytesIO(data))
                    images.append({"page": p_idx, "image": pil_img})
                except Exception as e:  # pylint: disable=broad-except
                    logger.error(
                        "[extract_images_from_pdf_bytes] image open error page=%s idx=%s err=%s",
                        p_idx,
                        i_idx,
                        e,
                    )
        except Exception as e:  # pylint: disable=broad-except
            logger.error("[extract_images_from_pdf_bytes] page=%s error=%s", p_idx, e)
    logger.info("[extract_images_from_pdf_bytes] end images_count=%s", len(images))
    return images


# ---------------- DOCX -----------------
def extract_text_from_docx_bytes(docx_bytes: bytes) -> List[Dict[str, str]]:
    logger.info("[extract_text_from_docx_bytes] start bytes_len=%s", len(docx_bytes or b""))
    doc = Document(io.BytesIO(docx_bytes))
    para_count = len(doc.paragraphs)
    logger.debug("[extract_text_from_docx_bytes] paragraphs=%s", para_count)
    sections: List[Dict[str, str]] = []
    current_name = "section:1"
    current_text = []
    index = 1
    for p_idx, para in enumerate(doc.paragraphs, start=1):
        is_heading = para.style.name.startswith("Heading") if para.style and para.style.name else False
        logger.debug("[extract_text_from_docx_bytes] para=%s is_heading=%s text_len=%s preview=%r",
                     p_idx, is_heading, len(para.text or ""), (para.text or "")[:200])
        if is_heading and para.text.strip():
            if current_text:
                section_text = "\n".join(current_text)
                sections.append({"source": current_name, "text": section_text})
                logger.debug("[extract_text_from_docx_bytes] new_section appended name=%s text_len=%s", current_name, len(section_text))
                current_text = []
                index += 1
            current_name = f"section:{para.text.strip()}"
        else:
            current_text.append(para.text)
    if current_text:
        section_text = "\n".join(current_text)
        sections.append({"source": current_name, "text": section_text})
        logger.debug("[extract_text_from_docx_bytes] final_section appended name=%s text_len=%s", current_name, len(section_text))
    if not sections:
        sections.append({"source": "section:1", "text": ""})
        logger.debug("[extract_text_from_docx_bytes] no sections -> appended empty section:1")
    logger.info("[extract_text_from_docx_bytes] end sections_count=%s", len(sections))
    return sections


def extract_images_from_docx(docx_bytes: bytes) -> List[Image.Image]:
    logger.info("[extract_images_from_docx] start bytes_len=%s", len(docx_bytes or b""))
    doc = Document(io.BytesIO(docx_bytes))
    rels_count = len(doc.part.rels.values())
    logger.debug("[extract_images_from_docx] rels_count=%s", rels_count)
    images: List[Image.Image] = []
    for r_idx, rel in enumerate(doc.part.rels.values(), start=1):
        target_ref = getattr(rel, "target_ref", "")
        has_image = "image" in str(target_ref)
        logger.debug("[extract_images_from_docx] rel_idx=%s target_ref=%s has_image=%s", r_idx, target_ref, has_image)
        if has_image:
            img_bytes = rel.target_part.blob
            logger.debug("[extract_images_from_docx] opening image blob_len=%s", len(img_bytes or b""))
            images.append(Image.open(io.BytesIO(img_bytes)))
    logger.info("[extract_images_from_docx] end images_count=%s", len(images))
    return images


# ---------------- PPTX -----------------
def extract_text_from_pptx_bytes(pptx_bytes: bytes) -> List[Dict[str, str]]:
    logger.info("[extract_text_from_pptx_bytes] start bytes_len=%s", len(pptx_bytes or b""))
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_total = len(prs.slides)
    logger.debug("[extract_text_from_pptx_bytes] slide_total=%s", slide_total)
    slides: List[Dict[str, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        shape_count = len(slide.shapes)
        logger.debug("[extract_text_from_pptx_bytes] slide=%s shape_count=%s has_notes=%s", i, shape_count, getattr(slide, "has_notes_slide", False))
        for s_idx, shape in enumerate(slide.shapes, start=1):
            has_tf = getattr(shape, "has_text_frame", False)
            if has_tf:
                texts.append(shape.text)
            logger.debug("[extract_text_from_pptx_bytes] slide=%s shape=%s has_text_frame=%s text_len=%s",
                         i, s_idx, has_tf, len(getattr(shape, "text", "") or ""))
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
            notes_txt = slide.notes_slide.notes_text_frame.text
            texts.append(notes_txt)
            logger.debug("[extract_text_from_pptx_bytes] slide=%s notes_text_len=%s", i, len(notes_txt or ""))
        joined = "\n".join(texts)
        logger.debug("[extract_text_from_pptx_bytes] slide=%s collected_text_len=%s preview=%r", i, len(joined), joined[:200])
        slides.append({"source": f"slide:{i}", "text": joined})
    logger.info("[extract_text_from_pptx_bytes] end slides_count=%s", len(slides))
    return slides


def extract_images_from_pptx(pptx_bytes: bytes) -> List[Dict[str, Image.Image]]:
    logger.info("[extract_images_from_pptx] start bytes_len=%s", len(pptx_bytes or b""))
    prs = Presentation(io.BytesIO(pptx_bytes))
    images: List[Dict[str, Image.Image]] = []
    for i, slide in enumerate(prs.slides, start=1):
        shape_count = len(slide.shapes)
        logger.debug("[extract_images_from_pptx] slide=%s shape_count=%s", i, shape_count)
        for s_idx, shape in enumerate(slide.shapes, start=1):
            is_picture = (shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
            logger.debug("[extract_images_from_pptx] slide=%s shape=%s is_picture=%s", i, s_idx, is_picture)
            if is_picture:
                blob_len = len(shape.image.blob or b"")
                logger.debug("[extract_images_from_pptx] slide=%s shape=%s image_blob_len=%s", i, s_idx, blob_len)
                images.append({"slide": i, "image": Image.open(io.BytesIO(shape.image.blob))})
    logger.info("[extract_images_from_pptx] end images_count=%s", len(images))
    return images


# ---------------- Legacy convert -----------------
def convert_legacy_office_to_modern(tmp_path: str) -> str:
    """LibreOffice ile eski ofis dosyasını modern formata çevirir."""
    logger.info("[convert_legacy_office_to_modern] start tmp_path=%s", tmp_path)
    outdir = os.path.dirname(tmp_path)
    cmd = ["libreoffice", "--headless", "--convert-to", "pdf", tmp_path, "--outdir", outdir]
    logger.debug("[convert_legacy_office_to_modern] cmd=%s", " ".join(cmd))
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    logger.debug("[convert_legacy_office_to_modern] returncode=%s stdout_len=%s stderr_len=%s",
                 result.returncode, len(result.stdout or b""), len(result.stderr or b""))
    if result.returncode != 0:
        err = result.stderr.decode()
        logger.error("[convert_legacy_office_to_modern] failed stderr=%s", err)
        raise RuntimeError(err)
    base, _ = os.path.splitext(tmp_path)
    out_path = base + ".pdf"
    logger.info("[convert_legacy_office_to_modern] end out_path=%s", out_path)
    return out_path


# ---------------- Ortak -----------------
def normalize_text(s: str) -> str:
    logger.debug("[normalize_text] start len=%s preview_before=%r", len(s or ""), (s or "")[:200])
    s = s.replace("\r", "\n")
    s = re.sub(r"[\t\v\f]+", " ", s)
    s = re.sub(r"\s+", " ", s)
    s = "".join(ch for ch in s if ch.isprintable())
    out = s.strip()
    logger.debug("[normalize_text] end len=%s preview_after=%r", len(out), out[:200])
    return out


def split_text_by_size(s: str, max_chars: int) -> List[str]:
    logger.debug("[split_text_by_size] start len=%s max_chars=%s", len(s or ""), max_chars)
    parts = [s[i:i + max_chars] for i in range(0, len(s), max_chars) if s[i:i + max_chars]]
    logger.debug("[split_text_by_size] end parts=%s", len(parts))
    return parts


def word_count(s: str) -> int:
    cnt = len(re.findall(r"\w+", s))
    logger.debug("[word_count] len=%s count=%s", len(s or ""), cnt)
    return cnt


def char_count(s: str) -> int:
    cnt = len(s)
    logger.debug("[char_count] len=%s", cnt)
    return cnt
