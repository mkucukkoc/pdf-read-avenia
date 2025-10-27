import os
import tempfile
import uuid
from datetime import datetime
import requests
from fastapi import HTTPException
from main import app, client, DEFAULT_MODEL, storage, DocRequest, parse_ppt_prompt, generate_random_style
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


@app.post("/generate-ppt")
async def generate_ppt(data: DocRequest):
    print("[/generate-ppt] 🌟 Sunum isteği alındı.")
    try:
        completion = client.chat.completions.create(
            model=DEFAULT_MODEL,
            messages=[
                {"role": "system", "content": """Sen bir sunum üreticisisin. Her slaytı şu formatta ver:
# Slide X
Title: ...
Content: ...
Image: (Bu başlıkla ilgili kısa bir sahne betimlemesi örn: \"kitap okuyan bir kadın\")"""},
                {"role": "user", "content": data.prompt}
            ],
            max_tokens=1500
        )
        generated_text = completion.choices[0].message.content.strip()
        slides = parse_ppt_prompt(generated_text)

        prs = Presentation()
        splash = prs.slides.add_slide(prs.slide_layouts[0])
        splash.shapes.title.text = f"📊 {data.prompt[:60]}..."
        splash.placeholders[1].text = "Bu sunum Avenia tarafından otomatik üretildi."

        for i, slide in enumerate(slides):
            print(f"[/generate-ppt] 📄 Slayt {i+1}: {slide['title'][:50]}...")
            s = prs.slides.add_slide(prs.slide_layouts[6])

            style = generate_random_style()  # prompt'tan tema yoksa rastgele seç

            fill = s.background.fill
            fill.solid()
            fill.fore_color.rgb = style["bg_color"]

            title_box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
            tf = title_box.text_frame
            tf.text = slide['title']
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.name = style["title_font"]
            tf.paragraphs[0].font.color.rgb = style["title_color"]

            date_box = s.shapes.add_textbox(Inches(8), Inches(0.1), Inches(2), Inches(0.3))
            dtf = date_box.text_frame
            dtf.text = datetime.now().strftime("%d %B %Y")
            dtf.paragraphs[0].font.size = Pt(12)
            dtf.paragraphs[0].font.name = 'Calibri'
            dtf.paragraphs[0].font.color.rgb = RGBColor(160, 160, 160)

            logo_path = "avenia_logo.png"
            if os.path.exists(logo_path):
                try:
                    s.shapes.add_picture(logo_path, Inches(0.1), Inches(5.3), height=Inches(0.5))
                except Exception as e:
                    print(f"[/generate-ppt] ⚠️ Logo eklenemedi: {e}")

            if slide['content']:
                content_box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4))
                ctf = content_box.text_frame
                ctf.text = slide['content']
                for p in ctf.paragraphs:
                    p.font.size = style["content_font_size"]
                    p.font.name = style["content_font"]
                    p.font.color.rgb = RGBColor(80, 80, 80)

            if slide['image']:
                try:
                    dalle_response = client.images.generate(
                        model="dall-e-3",
                        prompt=slide['image'],
                        n=1,
                        size="1024x1024",
                    )
                    image_url = dalle_response.data[0].url
                    image_data = requests.get(image_url).content
                    image_path = os.path.join(tempfile.gettempdir(), f"image_{uuid.uuid4().hex}.png")
                    with open(image_path, "wb") as f:
                        f.write(image_data)
                    s.shapes.add_picture(image_path, Inches(5.2), Inches(1.5), height=Inches(3.5))
                except Exception as e:
                    print("[/generate-ppt] ❌ Görsel oluşturulamadı:", str(e))

        filename = f"generated_{uuid.uuid4().hex}.pptx"
        filepath = os.path.join(tempfile.gettempdir(), filename)
        prs.save(filepath)

        bucket = storage.bucket()
        blob = bucket.blob(f"generated_ppts/{filename}")
        blob.upload_from_filename(filepath)
        blob.make_public()

        return {"status": "success", "file_url": blob.public_url}

    except Exception as e:
        print("[/generate-ppt] ❌ Hata oluştu:", str(e))
        raise HTTPException(status_code=500, detail=str(e))
