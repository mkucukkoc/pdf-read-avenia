import os
import io
import uuid
import tempfile
from fastapi import UploadFile, HTTPException
from main import app, storage
from pptx import Presentation
from fpdf import FPDF


@app.post("/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile):
    """Convert an uploaded PPTX file to PDF and return a Firebase URL."""
    try:
        ppt_bytes = await file.read()
        presentation = Presentation(io.BytesIO(ppt_bytes))

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)

        for idx, slide in enumerate(presentation.slides, start=1):
            pdf.add_page()
            pdf.multi_cell(0, 10, f"Slide {idx}")
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
            pdf.multi_cell(0, 10, "\n".join(texts))

        temp_dir = tempfile.gettempdir()
        filename = f"converted_{uuid.uuid4().hex}.pdf"
        filepath = os.path.join(temp_dir, filename)
        pdf.output(filepath)

        bucket = storage.bucket()
        blob = bucket.blob(f"converted_pdfs/{filename}")
        blob.upload_from_filename(filepath)
        blob.make_public()

        return {"status": "success", "file_url": blob.public_url}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
