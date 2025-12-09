import logging
import os
import io
import uuid
import tempfile
from fastapi import UploadFile, HTTPException
from main import app, storage
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger("pdf_read_refresh.endpoints.pdf_to_ppt")


@app.post("/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile):
    """Convert an uploaded PDF file to a PPTX presentation and return a Firebase URL."""
    logger.info(
        "PDF to PPT request",
        extra={"filename": getattr(file, "filename", None), "content_type": getattr(file, "content_type", None)},
    )
    try:
        pdf_bytes = await file.read()
        reader = PdfReader(io.BytesIO(pdf_bytes))

        prs = Presentation()
        layout = prs.slide_layouts[1]

        for idx, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            logger.debug("Page text extracted", extra={"page": idx, "length": len(text)})
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Page {idx}"
            body = slide.shapes.placeholders[1].text_frame
            first = True
            for line in text.split("\n"):
                cleaned = line.strip()
                if not cleaned:
                    continue
                if first:
                    body.text = cleaned
                    first = False
                else:
                    p = body.add_paragraph()
                    p.text = cleaned

        temp_dir = tempfile.gettempdir()
        filename = f"converted_{uuid.uuid4().hex}.pptx"
        filepath = os.path.join(temp_dir, filename)
        prs.save(filepath)
        logger.info("Presentation saved", extra={"path": filepath})

        bucket = storage.bucket()
        blob = bucket.blob(f"converted_ppts/{filename}")
        blob.upload_from_filename(filepath)
        blob.make_public()
        logger.info("Uploaded PPT to Firebase", extra={"url": blob.public_url})

        response = {"status": "success", "file_url": blob.public_url}
        logger.debug("PDF to PPT response", extra={"response": response})
        return response
    except Exception as e:
        logger.exception("PDF to PPT failed")
        raise HTTPException(status_code=500, detail=str(e))


