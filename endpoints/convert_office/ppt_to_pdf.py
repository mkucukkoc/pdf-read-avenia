import os
import io
import uuid
import tempfile
import logging
import hashlib  # sadece log için
from fastapi import UploadFile, HTTPException
from main import app, storage
from pptx import Presentation
from fpdf import FPDF

logger = logging.getLogger(__name__)


@app.post("/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile):
    """Convert an uploaded PPTX file to PDF and return a Firebase URL."""
    try:
        logger.info("[/ppt-to-pdf] START filename=%s content_type=%s",
                    getattr(file, "filename", None), getattr(file, "content_type", None))

        # 1) Dosyayı oku
        ppt_bytes = await file.read()
        logger.info("[/ppt-to-pdf] file.read ok bytes=%d md5=%s",
                    len(ppt_bytes), hashlib.md5(ppt_bytes).hexdigest())

        # 2) PPTX'i yükle
        presentation = Presentation(io.BytesIO(ppt_bytes))
        slide_count = len(presentation.slides)
        logger.info("[/ppt-to-pdf] PPTX parsed slide_count=%d", slide_count)

        # 3) PDF kurulum
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)
        logger.info("[/ppt-to-pdf] FPDF initialized (auto_page_break=True, margin=15, font=Arial:12)")

        # 4) Slaytları yaz
        for idx, slide in enumerate(presentation.slides, start=1):
            pdf.add_page()
            logger.info("[/ppt-to-pdf] slide=%d add_page()", idx)
            pdf.multi_cell(0, 10, f"Slide {idx}")

            texts = []
            shape_count = len(slide.shapes)
            logger.debug("[/ppt-to-pdf] slide=%d shapes=%d", idx, shape_count)

            for s_i, shape in enumerate(slide.shapes, start=1):
                has_text_attr = hasattr(shape, "text")
                logger.debug("[/ppt-to-pdf] slide=%d shape=%d has_text=%s", idx, s_i, has_text_attr)
                if has_text_attr:
                    t = shape.text
                    texts.append(t)
                    if t:
                        logger.debug("[/ppt-to-pdf] slide=%d shape=%d text_preview=%r", idx, s_i, t[:120])

            combined = "\n".join(texts)
            logger.debug("[/ppt-to-pdf] slide=%d combined_len=%d", idx, len(combined))
            pdf.multi_cell(0, 10, combined)

        # 5) Geçici dosyaya kaydet
        temp_dir = tempfile.gettempdir()
        filename = f"converted_{uuid.uuid4().hex}.pdf"
        filepath = os.path.join(temp_dir, filename)
        logger.info("[/ppt-to-pdf] saving PDF to temp path=%s", filepath)
        pdf.output(filepath)
        try:
            size_bytes = os.path.getsize(filepath)
        except OSError:
            size_bytes = -1
        logger.info("[/ppt-to-pdf] PDF saved size_bytes=%d", size_bytes)

        # 6) Firebase Storage'a yükle
        bucket = storage.bucket()
        logger.info("[/ppt-to-pdf] firebase bucket=%s", getattr(bucket, "name", "<unknown>"))
        blob_path = f"converted_pdfs/{filename}"
        blob = bucket.blob(blob_path)
        logger.info("[/ppt-to-pdf] uploading to storage path=%s", blob_path)
        blob.upload_from_filename(filepath)
        logger.info("[/ppt-to-pdf] upload complete")
        blob.make_public()
        logger.info("[/ppt-to-pdf] made public url=%s", blob.public_url)

        # 7) Yanıt
        logger.info("[/ppt-to-pdf] END success")
        return {"status": "success", "file_url": blob.public_url}

    except Exception as e:
        logger.exception("[/ppt-to-pdf] ERROR: %s", e)
        raise HTTPException(status_code=500, detail=str(e))










