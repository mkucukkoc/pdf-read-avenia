import os
import uuid
import aiohttp
from fastapi import HTTPException
from fastapi import Body
from main import app, client, DEFAULT_MODEL, save_embeddings_to_firebase
from pptx import Presentation
from firebase_admin import firestore


@app.post("/summarize-ppt-url/")
async def summarize_ppt_from_url(data: dict = Body(...)):
    print("🚀 [summarize-ppt-url] Endpoint tetiklendi")
    print("📦 Gelen data:", data)

    url = data.get("url")
    if not url:
        print("❌ URL bulunamadı!")
        raise HTTPException(status_code=400, detail="URL not provided")

    file_path = "temp.pptx"
    print(f"📥 Sunum indiriliyor: {url}")

    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url) as resp:
                if resp.status != 200:
                    print(f"❌ Dosya indirilemedi. HTTP status: {resp.status}")
                    raise HTTPException(status_code=500, detail="File download failed")
                with open(file_path, "wb") as f:
                    f.write(await resp.read())
        print("✅ Dosya indirildi:", file_path)
    except Exception as e:
        print("❌ Dosya indirme hatası:", e)
        raise HTTPException(status_code=500, detail="Download error")

    try:
        prs = Presentation(file_path)
        full_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + "\n"

        print("📝 Sunumdan çıkarılan içerik (ilk 200 karakter):")
        print(full_text[:200] or "[boş]")

        os.remove(file_path)
    except Exception as e:
        print("❌ PPTX okuma hatası:", e)
        raise HTTPException(status_code=500, detail="PowerPoint parse error")

    try:
        print("🤖 GPT-4 ile özetleniyor...")
        response = client.chat.completions.create(
            model=DEFAULT_MODEL,
            messages=[
                {"role": "system", "content": "Bu PowerPoint sunumunun içeriğini özetle:"},
                {"role": "user", "content": full_text[:3000]}
            ]
        )
        summary = response.choices[0].message.content
        print("✅ GPT özeti başarıyla alındı (ilk 200 karakter):")
        print(summary[:200])
        file_id = str(uuid.uuid4())
        user_id = data.get("user_id")
        chat_id = data.get("chat_id")
        print(chat_id,"chat_id")
        print(user_id,"user_id")
        db = firestore.client()
        save_embeddings_to_firebase(user_id, chat_id, file_id, full_text, summary, "PPTX")
        messages_ref = db.collection("users").document(user_id).collection("chats").document(chat_id).collection("messages")
        chat_ref = db.collection("users").document(user_id).collection("chats").document(chat_id)
        chat_ref.update({"file_id": file_id})
        print(f"[summarize-ppt-url] ✅ Chat doc file_id güncellendi: {file_id}")
        print(messages_ref,"messages_ref")
        messages_ref.add({
            "role": "assistant",
            "content": summary,  # GPT özeti
            "file_id": file_id,
            "timestamp": firestore.SERVER_TIMESTAMP
        })
        print(messages_ref,"messages_ref")
        return {"full_text": summary}
    except Exception as e:
        print("❌ GPT özetleme hatası:", e)
        raise HTTPException(status_code=500, detail="GPT summarization error")
