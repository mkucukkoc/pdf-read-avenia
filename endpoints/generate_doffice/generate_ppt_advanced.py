# endpoints/generate_ppt_advanced.py
import logging
import os
import re
import io
import json
import uuid
import math
import tempfile
from typing import List, Optional, Dict, Literal

import requests
from fastapi import HTTPException
from pydantic import BaseModel, Field, validator

from main import app, client, DEFAULT_MODEL, storage  # projendeki mevcut import yapısı
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

logger = logging.getLogger("pdf_read_refresh.endpoints.generate_ppt_advanced")


# ───────────────────────────────────────────────────────────────────────────────
# Pydantic Request & Plan Şeması
# ───────────────────────────────────────────────────────────────────────────────

class BrandKit(BaseModel):
    primary: Optional[str] = None
    secondary: Optional[str] = None
    accent: Optional[str] = None
    title_font: Optional[str] = None
    body_font: Optional[str] = None
    logo_url: Optional[str] = None

class PPTTheme(BaseModel):
    mode: Literal["light", "dark"] = "light"
    primary: str = "#1F6FEB"
    secondary: str = "#111827"
    accent: str = "#FDBA74"
    title_font: str = "Inter"
    body_font: str = "Inter"

class SlideSpec(BaseModel):
    type: Literal["title", "section", "bullets", "two-column", "image-focus",
                  "quote", "chart", "comparison", "timeline"] = "bullets"
    title: str = Field(..., description="Slide title (concise)")
    bullets: Optional[List[str]] = Field(default=None, description="List of bullet points")
    image_prompt: Optional[str] = Field(default=None, description="Image prompt text")
    chart_spec: Optional[Dict] = Field(default=None, description="Chart spec JSON")
    notes: Optional[str] = Field(default=None, description="Speaker notes for this slide")

class Plan(BaseModel):
    title: Optional[str] = None
    slides: List[SlideSpec]

class PPTAdvancedRequest(BaseModel):
    # içerik
    prompt: str
    language: Literal["tr", "en"] = "tr"
    audience: Optional[str] = None          # hedef kitle
    purpose: Optional[str] = None           # eğitim / pitch / iç iletişim vb.
    title: Optional[str] = None
    outline: Optional[List[str]] = None

    # üretim hedefleri
    slide_goal: int = 12
    charts_allowed: bool = True
    image_policy: Literal["generate", "none"] = "generate"
    image_style: Optional[str] = "clean, modern, flat illustration"
    speaker_notes: bool = True

    # mizanpaj / tema
    aspect_ratio: Literal["16:9", "4:3"] = "16:9"
    include_cover: bool = True
    include_agenda: bool = True
    include_summary: bool = True
    include_qna: bool = True
    include_closing: bool = True
    slide_numbers: bool = True
    header_text: Optional[str] = None
    footer_text: Optional[str] = None
    logo_url: Optional[str] = None
    theme: Optional[PPTTheme] = PPTTheme()
    brand_kit: Optional[BrandKit] = None

    # referanslar
    references: Optional[List[str]] = None

    @validator("slide_goal")
    def _slide_goal_min(cls, v):
        if v < 1:
            raise ValueError("slide_goal must be >= 1")
        return v


# ───────────────────────────────────────────────────────────────────────────────
# Yardımcılar
# ───────────────────────────────────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = (hex_color or "#000000").lstrip("#")
    if len(hex_color) < 6:
        hex_color = hex_color.zfill(6)
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)

def _apply_brand_kit(theme: PPTTheme, brand: Optional[BrandKit]) -> PPTTheme:
    if not brand:
        return theme
    # brand kit varsa temayı override et
    theme.primary   = brand.primary   or theme.primary
    theme.secondary = brand.secondary or theme.secondary
    theme.accent    = brand.accent    or theme.accent
    theme.title_font= brand.title_font or theme.title_font
    theme.body_font = brand.body_font  or theme.body_font
    return theme

def _new_presentation(aspect_ratio: str) -> Presentation:
    prs = Presentation()
    if aspect_ratio == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        # 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    return prs

def _download_to_temp(url: str, suffix: str) -> Optional[str]:
    try:
        r = requests.get(url, timeout=25)
        r.raise_for_status()
        path = os.path.join(tempfile.gettempdir(), f"dl_{uuid.uuid4().hex}{suffix}")
        with open(path, "wb") as f:
            f.write(r.content)
        return path
    except Exception as e:
        print("[/generate-ppt-advanced] ⚠️ Download failed:", e)
        return None

def _add_header_footer_boxes(slide, header_text: Optional[str], footer_text: Optional[str],
                             slide_no: Optional[int], theme: PPTTheme):
    # PowerPoint header/footer API sınırlı → textbox ile
    if header_text:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.5))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = header_text
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.size = Pt(11)
        p.font.name = theme.body_font
        p.font.color.rgb = _hex_to_rgb("#666666")

    footer_str = footer_text or ""
    if slide_no is not None:
        footer_str = (footer_str + "   ") if footer_str else ""
        footer_str += f"{slide_no}"

    if footer_str:
        fb = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.0), Inches(0.4))
        ft = fb.text_frame
        ft.clear()
        p2 = ft.paragraphs[0]
        p2.text = footer_str
        p2.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
        p2.font.size = Pt(10)
        p2.font.name = theme.body_font
        p2.font.color.rgb = _hex_to_rgb("#9CA3AF")

def _add_image(slide, image_path: str, left_in=7.2, top_in=1.5, height_in=3.5):
    try:
        slide.shapes.add_picture(image_path, Inches(left_in), Inches(top_in), height=Inches(height_in))
    except Exception as e:
        print("[/generate-ppt-advanced] ⚠️ Image add failed:", e)

def _generate_image(prompt: str) -> Optional[str]:
    try:
        dalle = client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            n=1,
            size="1024x1024",
        )
        image_url = dalle.data[0].url
        return _download_to_temp(image_url, ".png")
    except Exception as e:
        print("[/generate-ppt-advanced] ❌ Image generation failed:", e)
        return None

def _add_bullets_textbox(slide, text: List[str], theme: PPTTheme,
                         left=0.8, top=1.6, width=6.0, height=4.5, font_size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    for i, line in enumerate(text or []):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.name = theme.body_font
        p.font.size = Pt(font_size)
        p.font.color.rgb = _hex_to_rgb("#374151")

def _add_title(slide, title: str, theme: PPTTheme):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(10.5), Inches(1.0))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.name = theme.title_font
    p.font.color.rgb = _hex_to_rgb("#FFFFFF" if theme.mode == "dark" else "#111827")

def _add_subtitle(slide, subtitle: str, theme: PPTTheme):
    sb = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(10.8), Inches(1.0))
    tf = sb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.name = theme.body_font
    p.font.color.rgb = _hex_to_rgb("#9CA3AF")

def _add_quote(slide, text: str, theme: PPTTheme):
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(10.5), Inches(3.5))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"“{text}”"
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.size = Pt(28)
    p.font.name = theme.title_font
    p.font.color.rgb = _hex_to_rgb(theme.primary)
    p.font.italic = True

def _add_chart(slide, chart_spec: Dict):
    """
    chart_spec örneği:
    {
      "type": "bar",  # "line", "pie"
      "categories": ["Q1","Q2","Q3","Q4"],
      "series": [
        {"name":"Gelir","values":[10,12,15,13]},
        {"name":"Gider","values":[7,9,11,9]}
      ]
    }
    """
    try:
        ctype = chart_spec.get("type", "bar")
        categories = chart_spec.get("categories", [])
        series = chart_spec.get("series", [])

        chart_data = ChartData()
        chart_data.categories = categories

        for s in series:
            chart_data.add_series(s.get("name", ""), s.get("values", []))

        if ctype == "line":
            xlt = XL_CHART_TYPE.LINE_MARKERS
        elif ctype == "pie":
            xlt = XL_CHART_TYPE.PIE
        else:
            xlt = XL_CHART_TYPE.COLUMN_CLUSTERED

        x, y, cx, cy = Inches(5.8), Inches(1.8), Inches(6.5), Inches(4.5)
        chart = slide.shapes.add_chart(xlt, x, y, cx, cy, chart_data).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    except Exception as e:
        print("[/generate-ppt-advanced] ⚠️ Chart add failed:", e)

def _cover_slide(prs: Presentation, title: str, subtitle: str, theme: PPTTheme, logo_url: Optional[str]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # arka plan
    try:
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb("#0B1220" if theme.mode == "dark" else "#FFFFFF")
    except Exception:
        pass

    _add_title(s, title, theme)
    _add_subtitle(s, subtitle, theme)

    if logo_url:
        path = _download_to_temp(logo_url, ".png")
        if path:
            _add_image(s, path, left_in=0.6, top_in=6.2, height_in=0.6)

def _agenda_slide(prs: Presentation, plan: Plan, theme: PPTTheme, language="tr"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Agenda" if language == "en" else "Gündem", theme)
    titles = [sl.title for sl in plan.slides]
    _add_bullets_textbox(s, titles, theme)
    return s

def _summary_slide(prs: Presentation, plan: Plan, theme: PPTTheme, language="tr"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Key Takeaways" if language == "en" else "Öne Çıkanlar", theme)
    bullets = []
    for sl in plan.slides[:min(6, len(plan.slides))]:
        bullets.append(f"{sl.title} – " + (sl.bullets[0] if sl.bullets else ""))
    _add_bullets_textbox(s, bullets, theme)
    return s

def _qna_slide(prs: Presentation, theme: PPTTheme, language="tr"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Q&A" if language == "en" else "Soru & Cevap", theme)
    _add_bullets_textbox(s, [
        "—", "—", "—"
    ], theme)
    return s

def _closing_slide(prs: Presentation, theme: PPTTheme, language="tr", cta: Optional[str] = None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Thank You" if language == "en" else "Teşekkürler", theme)
    if cta:
        _add_bullets_textbox(s, [cta], theme)
    return s

def _render_content_slide(prs: Presentation, spec: SlideSpec, idx: int, theme: PPTTheme,
                          image_policy: str, speaker_notes: bool,
                          header_text: Optional[str], footer_text: Optional[str], slide_numbers: bool):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, spec.title, theme)

    # içerik
    if spec.type in ["bullets", "two-column", "section", "title", "comparison", "timeline"]:
        bullets = spec.bullets or []
        if spec.type == "two-column" and len(bullets) >= 4:
            mid = math.ceil(len(bullets)/2)
            left_bullets = bullets[:mid]
            right_bullets = bullets[mid:]
            _add_bullets_textbox(s, left_bullets, theme, left=0.8, top=1.8, width=5.5, height=4.5)
            _add_bullets_textbox(s, right_bullets, theme, left=6.2, top=1.8, width=5.5, height=4.5)
        else:
            _add_bullets_textbox(s, bullets, theme)

    if spec.type == "quote" and spec.bullets:
        _add_quote(s, " ".join(spec.bullets), theme)

    if spec.type == "chart" and spec.chart_spec:
        _add_chart(s, spec.chart_spec)

    # görsel
    if image_policy == "generate" and spec.image_prompt:
        img_path = _generate_image(spec.image_prompt)
        if img_path:
            _add_image(s, img_path)

    # speaker notes
    if speaker_notes and spec.notes:
        try:
            s.notes_slide.notes_text_frame.text = spec.notes
        except Exception:
            pass

    # header/footer
    _add_header_footer_boxes(
        s,
        header_text=header_text,
        footer_text=footer_text,
        slide_no=(idx if slide_numbers else None),
        theme=theme
    )
    return s


def _extract_json(text: str) -> str:
    # Model "sadece JSON" desek de bazen uymayabilir. İlk JSON bloğunu çek.
    try:
        json.loads(text)
        return text
    except Exception:
        pass
    m = re.search(r'\{[\s\S]*\}', text)
    if m:
        return m.group(0)
    raise ValueError("JSON plan parse edilemedi.")

def _ensure_slide_goal(plan: Plan, goal: int) -> Plan:
    slides = plan.slides[:]
    # azsa: çok bullet'lı slaytları böl
    while len(slides) < goal:
        if not slides:
            break
        idx = max(range(len(slides)), key=lambda i: len(slides[i].bullets or []))
        cand = slides[idx]
        if not cand.bullets or len(cand.bullets) < 4:
            break
        mid = len(cand.bullets)//2
        left = SlideSpec(type=cand.type, title=cand.title + " (1)", bullets=cand.bullets[:mid],
                         image_prompt=cand.image_prompt, notes=cand.notes, chart_spec=cand.chart_spec)
        right = SlideSpec(type=cand.type, title=cand.title + " (2)", bullets=cand.bullets[mid:],
                          image_prompt=cand.image_prompt, notes=cand.notes, chart_spec=cand.chart_spec)
        slides.pop(idx)
        slides.insert(idx, right)
        slides.insert(idx, left)
    # fazlaysa: sondan birleştir
    while len(slides) > goal and len(slides) >= 2:
        a = slides[-2]
        b = slides[-1]
        merged_bullets = (a.bullets or []) + (b.bullets or [])
        a = SlideSpec(
            type=a.type if a.type != "title" else "bullets",
            title=a.title,
            bullets=merged_bullets[:10],
            image_prompt=a.image_prompt or b.image_prompt,
            notes=((a.notes or "") + ("\n" + (b.notes or "") if b.notes else "")) or None,
            chart_spec=a.chart_spec or b.chart_spec
        )
        slides = slides[:-2] + [a]
    return Plan(title=plan.title, slides=slides)

def _build_plan_prompt(data: PPTAdvancedRequest):
    lang = data.language
    goal = data.slide_goal
    outline_str = "\n".join([f"- {o}" for o in (data.outline or [])]) if data.outline else ""
    outline_block = outline_str if outline_str else "(If outline missing, propose a strong outline.)"

    # modelden sadece JSON isteme
    system = (
        "You are a world-class presentation designer. "
        "Return ONLY JSON matching the schema. No prose, no markdown, no comments."
    )
    user = f"""
LANGUAGE: {lang}
AUDIENCE: {data.audience or "-"}
PURPOSE: {data.purpose or "-"}
TITLE: {data.title or 'Presentation'}
SLIDE_GOAL: exactly {goal} slides
IMAGE_STYLE: {data.image_style}
CHARTS_ALLOWED: {str(data.charts_allowed).lower()}

OUTLINE:
{outline_block}

JSON SCHEMA (example):
{{
  "title": "string",
  "slides": [
    {{
      "type": "bullets|two-column|image-focus|section|quote|chart|title|comparison|timeline",
      "title": "string",
      "bullets": ["short bullet 1", "short bullet 2"],
      "image_prompt": "optional prompt",
      "chart_spec": {{
        "type": "bar|line|pie",
        "categories": ["A","B"],
        "series": [{{"name":"S1","values":[1,2]}}]
      }},
      "notes": "speaker notes"
    }}
  ]
}}

CONSTRAINTS:
- Exactly {goal} slides
- Titles concise; bullets <= 5 and <= 12 words each
- Use mixed slide types (sections, two-column, quotes, charts when helpful)
- Respect LANGUAGE for all text and notes
"""
    return system, user


# ───────────────────────────────────────────────────────────────────────────────
# Endpoint
# ───────────────────────────────────────────────────────────────────────────────

@app.post("/generate-ppt-advanced")
async def generate_ppt_advanced(data: PPTAdvancedRequest):
    logger.info(
        "Generate PPT advanced request received",
        extra={"prompt_preview": (data.prompt or "")[:200], "slide_goal": data.slide_goal},
    )
    warnings = []

    try:
        # 0) Tema + brand kit
        logger.info("Applying theme and brand kit")
        theme = data.theme or PPTTheme()
        theme = _apply_brand_kit(theme, data.brand_kit)
        if (not data.logo_url) and (data.brand_kit and data.brand_kit.logo_url):
            data.logo_url = data.brand_kit.logo_url

        # 1) Plan üret
        logger.info("Requesting plan generation")
        system, user = _build_plan_prompt(data)
        raw = client.chat.completions.create(
            model=DEFAULT_MODEL,
            messages=[{"role": "system", "content": system},
                      {"role": "user", "content": user}],
            temperature=0.6,
            max_completion_tokens=3500
        ).choices[0].message.content

        logger.debug("Plan raw length", extra={"length": len(raw or "")})
        try:
            plan_json = _extract_json(raw)
            plan_obj = Plan.parse_obj(json.loads(plan_json))
        except Exception as e:
            logger.exception("Failed to parse plan JSON")
            raise HTTPException(status_code=500, detail="Model planını JSON olarak parse edemedim.")

        # 1.b) Slide sayısını hedefe dengele
        if len(plan_obj.slides) != data.slide_goal:
            logger.warning(
                "Slide count mismatch - balancing",
                extra={"got": len(plan_obj.slides), "want": data.slide_goal},
            )
            plan_obj = _ensure_slide_goal(plan_obj, data.slide_goal)

        # 2) Sunumu oluştur
        logger.info("Building presentation from plan")
        prs = _new_presentation(data.aspect_ratio)

        # Kapak
        slide_idx = 1
        if data.include_cover:
            _ = _cover_slide(prs, data.title or "Presentation", data.prompt, theme, data.logo_url)
            slide_idx += 1

        # Agenda
        if data.include_agenda:
            _ = _agenda_slide(prs, plan_obj, theme, language=data.language)
            slide_idx += 1

        # İçerik slaytları
        for i, spec in enumerate(plan_obj.slides, start=1):
            _ = _render_content_slide(
                prs,
                spec,
                idx=slide_idx,
                theme=theme,
                image_policy=data.image_policy,
                speaker_notes=data.speaker_notes,
                header_text=data.header_text,
                footer_text=data.footer_text,
                slide_numbers=data.slide_numbers,
            )
            slide_idx += 1

        # Özet
        if data.include_summary:
            _ = _summary_slide(prs, plan_obj, theme, language=data.language)
            slide_idx += 1

        # Q&A
        if data.include_qna:
            _ = _qna_slide(prs, theme, language=data.language)
            slide_idx += 1

        # Kapanış
        if data.include_closing:
            _ = _closing_slide(prs, theme, language=data.language, cta=None)
            slide_idx += 1

        # 3) Geçici dosyaya kaydet
        temp_dir = tempfile.gettempdir()
        filename = f"generated_{uuid.uuid4().hex}.pptx"
        filepath = os.path.join(temp_dir, filename)
        prs.save(filepath)
        logger.info("PPTX saved", extra={"filepath": filepath})

        # 4) Firebase Storage’a yükle
        logger.info("Uploading PPTX to Firebase Storage")
        bucket = storage.bucket()
        blob = bucket.blob(f"generated_ppts/{filename}")
        blob.upload_from_filename(filepath)
        blob.make_public()
        logger.info("Firebase upload completed", extra={"file_url": blob.public_url})

        response_payload = {"status": "success", "file_url": blob.public_url, "warnings": warnings}
        logger.debug("Generate PPT advanced response", extra={"response": response_payload})
        return response_payload

    except Exception as e:
        logger.exception("Generate PPT advanced failed")
        raise HTTPException(status_code=500, detail=str(e))




