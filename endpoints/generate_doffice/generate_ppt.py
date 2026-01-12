import json
import logging
import os
import tempfile
import uuid
import time
from datetime import datetime
from typing import Any, Dict, List, Optional

import httpx
from fastapi import APIRouter, HTTPException, Request
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from firebase_admin import storage

from core.gemini_prompt import build_system_message, merge_parts_with_system
from core.firebase import db
from schemas import PptRequest
from endpoints.logging.utils_logging import log_gemini_request, log_gemini_response
from endpoints.generate_doffice.ppt_style import generate_random_style
from usage_tracking import build_base_event, finalize_event, parse_gemini_usage, enqueue_usage_update

logger = logging.getLogger("pdf_read_refresh.endpoints.generate_ppt")
router = APIRouter()

SYSTEM_INSTRUCTION = (
    "You are a professional presentation generator. "
    'Return ONLY valid JSON with this schema: {"slides": [{"title": "...", "bullets": ["..."], "image_prompt": "..."}]}. '
    "Each slide must have a short title, 3-5 concise bullet points, and an optional image_prompt for illustration. "
    "No markdown or extra text outside the JSON."
)


def _effective_model() -> str:
    model = os.getenv("GEMINI_PPT_MODEL") or os.getenv("GEMINI_SEARCH_MODEL") or "models/gemini-2.5-pro"
    if not model.startswith("models/"):
        model = f"models/{model}"
    return model


def _build_usage_context(request: Request, user_id: str, model: str) -> Optional[Dict[str, Any]]:
    if not user_id:
        return None
    token_payload = getattr(request.state, "token_payload", {}) or {}
    request_id = (
        request.headers.get("x-request-id")
        or request.headers.get("x-requestid")
        or f"req_{uuid.uuid4().hex}"
    )
    return build_base_event(
        request_id=request_id,
        user_id=user_id,
        endpoint="generate_ppt",
        provider="gemini",
        model=model,
        token_payload=token_payload,
        request=request,
    )


def _enqueue_usage_event(
    usage_context: Optional[Dict[str, Any]],
    usage_data: Dict[str, int],
    latency_ms: int,
    *,
    status: str,
    error_code: Optional[str],
) -> None:
    if not usage_context or not db:
        return
    try:
        event = finalize_event(
            usage_context,
            input_tokens=usage_data.get("inputTokens", 0),
            output_tokens=usage_data.get("outputTokens", 0),
            latency_ms=latency_ms,
            status=status,
            error_code=error_code,
        )
        enqueue_usage_update(db, event)
    except Exception:
        logger.warning("Usage tracking failed for generate_ppt", exc_info=True)


async def _call_gemini_json(prompt: str, system_message: Optional[str]) -> tuple[Dict[str, Any], Dict[str, Any]]:
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise HTTPException(
            status_code=500,
            detail={"success": False, "error": "gemini_api_key_missing", "message": "GEMINI_API_KEY is required"},
        )

    model = _effective_model()
    url = f"https://generativelanguage.googleapis.com/v1beta/{model}:generateContent?key={api_key}"
    parts = [
        {"text": SYSTEM_INSTRUCTION},
        {"text": f"USER_PROMPT: {prompt}"},
    ]
    effective_parts = merge_parts_with_system(parts, system_message)
    payload = {
        "contents": [
            {
                "role": "user",
                "parts": effective_parts,
            }
        ],
        "tools": [{"google_search": {}}],
        "generationConfig": {
            "responseMimeType": "application/json",
            "temperature": 0.5,
            "maxOutputTokens": 2048,
        },
    }

    log_gemini_request(
        logger,
        "generate_ppt",
        url=url,
        payload=payload,
        model=model,
    )
    logger.info("Gemini PPT JSON request", extra={"model": model, "prompt_preview": prompt[:200]})

    async with httpx.AsyncClient(timeout=90) as client:
        resp = await client.post(url, json=payload)

    body_preview = (resp.text or "")[:800]
    response_json = resp.json() if resp.text else {}
    log_gemini_response(
        logger,
        "generate_ppt",
        url=url,
        status_code=resp.status_code,
        response=response_json,
    )
    logger.info("Gemini PPT JSON response", extra={"status": resp.status_code, "body_preview": body_preview})

    if not resp.ok:
        raise HTTPException(
            status_code=resp.status_code,
            detail={"success": False, "error": "gemini_ppt_failed", "message": body_preview},
        )

    data = response_json
    parts = (data.get("candidates") or [{}])[0].get("content", {}).get("parts", [])
    if not parts:
        raise HTTPException(
            status_code=500,
            detail={"success": False, "error": "empty_response", "message": "Gemini returned no content"},
        )
    try:
        parsed = json.loads(parts[0].get("text") or "{}")
    except json.JSONDecodeError:
        raise HTTPException(
            status_code=500,
            detail={"success": False, "error": "invalid_json", "message": "Gemini returned non-JSON content"},
        )
    return parsed, response_json


def _build_ppt_from_json(ppt_data: Dict[str, Any], prompt: str) -> str:
    slides: List[Dict[str, Any]] = ppt_data.get("slides") or []
    prs = Presentation()

    splash = prs.slides.add_slide(prs.slide_layouts[0])
    splash.shapes.title.text = f"📊 {prompt[:60]}..."
    splash.placeholders[1].text = "Bu sunum Avenia tarafından otomatik üretildi."

    for i, slide in enumerate(slides):
        title = (slide or {}).get("title") or f"Slide {i+1}"
        bullets = (slide or {}).get("bullets") or []
        image_prompt = (slide or {}).get("image_prompt") or ""

        logger.debug("Adding slide", extra={"index": i + 1, "title": title[:50]})
        s = prs.slides.add_slide(prs.slide_layouts[6])
        style = generate_random_style()

        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = style["bg_color"]

        title_box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = style["title_font"]
        tf.paragraphs[0].font.color.rgb = style["title_color"]

        date_box = s.shapes.add_textbox(Inches(8), Inches(0.1), Inches(2), Inches(0.3))
        dtf = date_box.text_frame
        dtf.text = datetime.now().strftime("%d %B %Y")
        dtf.paragraphs[0].font.size = Pt(12)
        dtf.paragraphs[0].font.name = "Calibri"
        dtf.paragraphs[0].font.color.rgb = RGBColor(160, 160, 160)

        logo_path = "avenia_logo.png"
        if os.path.exists(logo_path):
            try:
                s.shapes.add_picture(logo_path, Inches(0.1), Inches(5.3), height=Inches(0.5))
            except Exception as e:
                logger.warning("Failed to add logo to slide", extra={"error": str(e)})

        if bullets:
            content_box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(4))
            ctf = content_box.text_frame
            ctf.text = ""
            for b in bullets:
                p = ctf.add_paragraph()
                p.text = str(b)
                p.font.size = style["content_font_size"]
                p.font.name = style["content_font"]
                p.font.color.rgb = RGBColor(80, 80, 80)

        if image_prompt:
            # Placeholder: görsel üretimi gereksinimi varsa burada ekleyebiliriz (şimdilik pasif)
            logger.debug("Image prompt captured for slide", extra={"index": i + 1, "image_prompt": image_prompt[:120]})

    filename = f"generated_{uuid.uuid4().hex}.pptx"
    filepath = os.path.join(tempfile.gettempdir(), filename)
    prs.save(filepath)
    return filepath, filename


@router.post("/generate-ppt")
async def generate_ppt(data: PptRequest, request: Request):
    logger.info("Generate PPT request received", extra={"prompt_length": len(data.prompt or "")})
    payload = getattr(request.state, "token_payload", {}) or {}
    user_id = payload.get("uid") or payload.get("userId") or payload.get("sub") or ""
    usage_context: Optional[Dict[str, Any]] = None
    usage_data: Dict[str, int] = {}
    start_time = time.monotonic()
    status = "success"
    error_code = None
    try:
        # 1) Gemini'den yapılandırılmış JSON al
        model = _effective_model()
        usage_context = _build_usage_context(request, user_id, model)
        system_message = build_system_message(
            language=None,
            tone_key=data.tone_key,
            response_style=None,
            include_response_style=False,
            include_followup=False,
        )
        ppt_json, response_json = await _call_gemini_json(data.prompt, system_message)
        usage_data = parse_gemini_usage(response_json)
        logger.debug("Gemini PPT JSON parsed", extra={"keys": list(ppt_json.keys())})

        # 2) PowerPoint dosyasını oluştur
        filepath, filename = _build_ppt_from_json(ppt_json, data.prompt)
        logger.info("PPT built", extra={"filepath": filepath})

        # 3) Firebase'e yükle
        bucket = storage.bucket()
        blob = bucket.blob(f"generated_ppts/{filename}")
        blob.upload_from_filename(filepath)
        blob.make_public()
        logger.info("Presentation uploaded to Firebase", extra={"file_url": blob.public_url})

        response_payload = {
            "status": "success",
            "file_url": blob.public_url,
            "slide_count": len(ppt_json.get("slides") or []),
        }
        logger.debug("Generate PPT response payload", extra={"response": response_payload})
        return response_payload

    except HTTPException as exc:
        status = "error"
        if isinstance(exc.detail, dict):
            error_code = exc.detail.get("error")
        error_code = error_code or "generate_ppt_failed"
        logger.exception("Generate PPT failed (HTTPException)")
        raise
    except Exception as e:
        status = "error"
        error_code = "generate_ppt_failed"
        logger.exception("Generate PPT failed")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        _enqueue_usage_event(
            usage_context,
            usage_data,
            int((time.monotonic() - start_time) * 1000),
            status=status,
            error_code=error_code,
        )


