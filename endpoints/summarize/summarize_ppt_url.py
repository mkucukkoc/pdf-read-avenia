import logging
import os
import uuid
import aiohttp
from fastapi import HTTPException
from fastapi import Body
from main import app, client, DEFAULT_MODEL, save_embeddings_to_firebase
from pptx import Presentation
from firebase_admin import firestore

logger = logging.getLogger("pdf_read_refresh.endpoints.summarize_ppt_url")


@app.post("/summarize-ppt-url/")
async def summarize_ppt_from_url(data: dict = Body(...)):
    logger.info("Summarize PPT URL request received", extra={"data": data})

    url = data.get("url")
    if not url:
        logger.warning("PPT summarize request missing URL")
        raise HTTPException(status_code=400, detail="URL not provided")

    file_path = "temp.pptx"
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url) as resp:
                if resp.status != 200:
                    logger.error("Failed to download PPT file", extra={"status": resp.status})
                    raise HTTPException(status_code=500, detail="File download failed")
                with open(file_path, "wb") as f:
                    f.write(await resp.read())
        logger.info("PPT file downloaded", extra={"file_path": file_path})
    except Exception as e:
        logger.exception("PPT file download error")
        raise HTTPException(status_code=500, detail="Download error")

    try:
        prs = Presentation(file_path)
        full_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + "\n"
        logger.debug("Extracted PPT text", extra={"length": len(full_text)})
        os.remove(file_path)
    except Exception as e:
        logger.exception("PPT parse error")
        raise HTTPException(status_code=500, detail="PowerPoint parse error")

    try:
        response = client.chat.completions.create(
            model=DEFAULT_MODEL,
            messages=[
                {"role": "system", "content": "Bu PowerPoint sunumunun içeriğini özetle:"},
                {"role": "user", "content": full_text[:3000]},
            ],
        )
        summary = response.choices[0].message.content
        logger.info("PPT summary generated", extra={"summary_length": len(summary)})
        file_id = str(uuid.uuid4())
        user_id = data.get("user_id")
        chat_id = data.get("chat_id")
        db = firestore.client()
        save_embeddings_to_firebase(user_id, chat_id, file_id, full_text, summary, "PPTX")
        messages_ref = (
            db.collection("users")
            .document(user_id)
            .collection("chats")
            .document(chat_id)
            .collection("messages")
        )
        chat_ref = db.collection("users").document(user_id).collection("chats").document(chat_id)
        chat_ref.update({"file_id": file_id})
        messages_ref.add(
            {
                "role": "assistant",
                "content": summary,
                "file_id": file_id,
                "timestamp": firestore.SERVER_TIMESTAMP,
            }
        )
        response_payload = {"full_text": summary}
        logger.debug("Summarize PPT response payload", extra={"response": response_payload})
        return response_payload
    except Exception as e:
        logger.exception("GPT summarization error for PPT")
        raise HTTPException(status_code=500, detail="GPT summarization error")









